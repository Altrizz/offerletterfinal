# app.py — Streamlit Cloud ready
# Template selector (with previews) • Custom upload
# Robust PPTX replacement (split runs/tables/groups)
# Session History with thumbnails • Export ZIP • Restore

import re
import io
import os
import base64
import zipfile
from pathlib import Path
from datetime import date, datetime
from io import BytesIO
from typing import List, Optional, Tuple

import streamlit as st
from pptx import Presentation
from pptx.shapes.group import GroupShape

# ==============================
# Formatting helpers
# ==============================
MESES_ES = [
    "enero","febrero","marzo","abril","mayo","junio",
    "julio","agosto","septiembre","octubre","noviembre","diciembre",
]
def fecha_es(d: date) -> str:
    return f"{d.day} de {MESES_ES[d.month-1]} de {d.year}"

def format_ars_dots(value) -> str:
    digits = "".join(ch for ch in str(value) if ch.isdigit())
    if not digits:
        return str(value)
    return f"{int(digits):,}".replace(",", ".")

# ==============================
# Placeholder patterns
# ==============================
PLACEHOLDER = re.compile(r"{{\s*([A-Z0-9_]+)\s*}}")
PAT_NAME = re.compile(r"\{X{6}\}")            # {XXXXXX}
PAT_POS  = re.compile(r"(?<!\{)X{8}(?!\})")   # XXXXXXXX (not inside {})
PAT_DATE = re.compile(r"\bX{2}\s+de\s+X{4,5}\s+de\s+\d{4}\b")
PAT_SAL  = re.compile(r"\bX\.XXX\.XXX\b")

def apply_x_style(text: str, mapping: dict) -> str:
    name = mapping.get("CANDIDATE_NAME")
    position = mapping.get("POSITION")
    join_date_es = mapping.get("JOIN_DATE")
    salary = mapping.get("SALARY")
    city = mapping.get("CITY", "Buenos Aires")
    offer_date_es = mapping.get("DATE")

    out = text
    if name:     out = PAT_NAME.sub(name, out)
    if position: out = PAT_POS.sub(position, out)
    if join_date_es: out = PAT_DATE.sub(join_date_es, out)
    if salary:   out = PAT_SAL.sub(format_ars_dots(salary), out)

    striped = out.strip()
    if striped == ", Buenos Aires" or striped.endswith(", Buenos Aires"):
        out = f"{offer_date_es}, {city}" if offer_date_es else city
    return out

def replace_placeholders_in_text(text: str, mapping: dict) -> str:
    def repl(m):
        key = m.group(1).upper()
        return str(mapping.get(key, m.group(0)))
    return apply_x_style(PLACEHOLDER.sub(repl, text), mapping)

# ==============================
# PPTX text replacement (runs/tables/groups safe)
# ==============================
def _replace_in_text_frame(tf, mapping: dict):
    for para in tf.paragraphs:
        if not para.runs:
            if para.text:
                new = replace_placeholders_in_text(para.text, mapping)
                if new != para.text:
                    para.text = new
            continue
        full = "".join(run.text for run in para.runs)
        new = replace_placeholders_in_text(full, mapping)
        if new == full:
            continue
        para.runs[0].text = new
        for r in para.runs[1:]:
            r.text = ""

def _replace_in_table(tbl, mapping: dict):
    for r in tbl.rows:
        for c in r.cells:
            if c.text_frame:
                _replace_in_text_frame(c.text_frame, mapping)

def _walk_shapes(shapes, mapping: dict):
    for shape in shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            _replace_in_text_frame(shape.text_frame, mapping)
        if getattr(shape, "has_table", False):
            _replace_in_table(shape.table, mapping)
        if isinstance(shape, GroupShape):
            _walk_shapes(shape.shapes, mapping)

def render_pptx(pptx_bytes: bytes, mapping: dict) -> BytesIO:
    prs = Presentation(BytesIO(pptx_bytes))
    for slide in prs.slides:
        _walk_shapes(slide.shapes, mapping)
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out

# ==============================
# Template loader & thumbnail utils
# ==============================
def list_templates() -> List[Path]:
    tpl_dir = Path(__file__).with_name("templates")
    if not tpl_dir.exists():
        return []
    return sorted([p for p in tpl_dir.glob("*.pptx") if p.is_file()])

def load_template_bytes(label: str) -> Optional[bytes]:
    for p in list_templates():
        if p.stem.replace("_", " ") == label:
            return p.read_bytes()
    return None

def pptx_doc_thumbnail(pptx_bytes: bytes) -> Optional[Tuple[str, bytes]]:
    """
    Extract document preview thumbnail from /docProps/thumbnail.jpeg|jpg|png if present.
    Returns (mime, data) or None.
    """
    try:
        with zipfile.ZipFile(BytesIO(pptx_bytes)) as zf:
            for name in ("docProps/thumbnail.jpeg", "docProps/thumbnail.jpg", "docProps/thumbnail.png"):
                if name in zf.namelist():
                    data = zf.read(name)
                    mime = "image/jpeg" if name.endswith(("jpeg", "jpg")) else "image/png"
                    return mime, data
    except Exception:
        pass
    return None

def svg_placeholder(title: str, subtitle: str = "", w: int = 480, h: int = 270) -> bytes:
    title = (title or "").replace("&", "&amp;")
    subtitle = (subtitle or "").replace("&", "&amp;")
    svg = f"""
    <svg xmlns="http://www.w3.org/2000/svg" width="{w}" height="{h}">
      <defs>
        <linearGradient id="g" x1="0" y1="0" x2="1" y2="1">
          <stop offset="0%" stop-color="#F2F4F8"/>
          <stop offset="100%" stop-color="#E9ECF3"/>
        </linearGradient>
      </defs>
      <rect width="100%" height="100%" rx="16" fill="url(#g)" stroke="#D9DFEA"/>
      <text x="24" y="120" font-family="Inter,system-ui,Segoe UI,Arial" font-weight="700"
            font-size="28" fill="#0b1220">{title}</text>
      <text x="24" y="160" font-family="Inter,system-ui,Segoe UI,Arial"
            font-size="16" fill="#475569">{subtitle}</text>
      <rect x="24" y="190" width="{w-48}" height="8" rx="4" fill="#E2E8F0"/>
      <rect x="24" y="210" width="{int((w-48)*0.65)}" height="8" rx="4" fill="#EEF2F7"/>
    </svg>
    """.strip()
    return svg.encode("utf-8")

def embed_svg(svg_bytes: bytes, width: int = 240):
    b64 = base64.b64encode(svg_bytes).decode("utf-8")
    st.markdown(f"<img src='data:image/svg+xml;base64,{b64}' width='{width}' />", unsafe_allow_html=True)

def first_texts_from_pptx(pptx_bytes: bytes, max_len: int = 40) -> Tuple[str, str]:
    """Best-effort: pull first two meaningful text lines from first slide."""
    try:
        prs = Presentation(BytesIO(pptx_bytes))
        if not prs.slides:
            return "", ""
        texts = []
        for shp in prs.slides[0].shapes:
            if getattr(shp, "has_text_frame", False):
                for p in shp.text_frame.paragraphs:
                    s = "".join(r.text for r in p.runs) if p.runs else (p.text or "")
                    s = " ".join(s.split())
                    if s:
                        texts.append(s)
            if len(texts) >= 2:
                break
        t1 = texts[0][:max_len] if texts else "Offer Template"
        t2 = texts[1][:max_len] if len(texts) > 1 else ""
        return t1, t2
    except Exception:
        return "Offer Template", ""

# ==============================
# Session History
# ==============================
HISTORY_KEY = "offer_history"
MAX_HISTORY = 50

def ensure_history():
    if HISTORY_KEY not in st.session_state:
        st.session_state[HISTORY_KEY] = []

def push_history(entry: dict):
    ensure_history()
    st.session_state[HISTORY_KEY].append(entry)
    if len(st.session_state[HISTORY_KEY]) > MAX_HISTORY:
        st.session_state[HISTORY_KEY] = st.session_state[HISTORY_KEY][-MAX_HISTORY:]

def delete_history(idx: int):
    ensure_history()
    if 0 <= idx < len(st.session_state[HISTORY_KEY]):
        st.session_state[HISTORY_KEY].pop(idx)

def zip_all_history() -> bytes:
    ensure_history()
    mem = io.BytesIO()
    with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as zf:
        for item in st.session_state[HISTORY_KEY]:
            zf.writestr(item["file_name"], item["pptx_bytes"])
    mem.seek(0)
    return mem.getvalue()

# ==============================
# UI
# ==============================
st.set_page_config(page_title="Offer Letter Generator", page_icon="📄", layout="centered")
st.title("Offer Letter Generator")

# --- Template picker with previews ---
templates = list_templates()
tpl_labels = [p.stem.replace("_", " ") for p in templates]
choices = ["— Select built-in template —"] + tpl_labels + ["Upload custom…"]
pick = st.selectbox("Offer letter template", options=choices, index=0)

uploaded_file = None
builtin_template_bytes: Optional[bytes] = None

if pick == "Upload custom…":
    uploaded_file = st.file_uploader(
        "Upload a PPTX template",
        type=["pptx"],
        help="Custom .pptx with placeholders like {{CANDIDATE_NAME}}, {{POSITION}}, etc."
    )
elif pick != "— Select built-in template —":
    builtin_template_bytes = load_template_bytes(pick)
    if builtin_template_bytes is None:
        st.error("Template not found or unreadable. Check the templates/ folder.")
    else:
        st.success(f"Using built-in template: **{pick}**")

# Preview gallery (built-in templates)
if templates:
    st.caption("Template previews")
    cols = st.columns(3)
    for i, p in enumerate(templates):
        with cols[i % 3]:
            pptx_bytes = p.read_bytes()
            thumb = pptx_doc_thumbnail(pptx_bytes)
            if thumb:
                mime, data = thumb
                # Use new parameter name to avoid deprecation warning
                st.image(data, caption=p.stem.replace("_", " "), use_container_width=True)
            else:
                t1, t2 = first_texts_from_pptx(pptx_bytes)
                embed_svg(svg_placeholder(t1, t2), width=300)
                st.caption(p.stem.replace("_", " "))

# --- Inputs ---
colA, colB = st.columns(2)
with colA:
    first_name = st.text_input("First name", placeholder="Jane")
    position   = st.text_input("Position", placeholder="Software Engineer")
    salary_num = st.number_input("Salary (ARS)", min_value=0, value=2000000, step=5000)
with colB:
    last_name  = st.text_input("Last name", placeholder="Doe")
    offer_date = st.date_input("Offer date", value=date.today())
    join_date  = st.date_input("Join date",  value=date.today())

common_cities = ["Buenos Aires", "Córdoba", "Rosario", "Mendoza", "Other..."]
city_pick = st.selectbox("City", options=common_cities, index=0)
city = st.text_input("City (custom)", value="", placeholder="Type city") if city_pick == "Other..." else city_pick

st.subheader("Extra placeholders (optional)")
extras = st.data_editor([{"key": "", "value": ""}], num_rows="dynamic", hide_index=True)

# --- Actions ---
col1, col2 = st.columns(2)
with col1:
    generate_clicked = st.button(
        "Generate Offer Letter",
        type="primary",
        disabled=(not builtin_template_bytes and not uploaded_file)
    )
with col2:
    if st.button("Clear fields"):
        st.experimental_rerun()

# --- Generate ---
if generate_clicked:
    # Source PPTX bytes
    if builtin_template_bytes:
        source_bytes = builtin_template_bytes
        template_label = pick
    elif uploaded_file:
        source_bytes = uploaded_file.read()
        template_label = uploaded_file.name
    else:
        st.warning("Please select a built-in template or upload a PPTX.")
        st.stop()

    full_name = f"{first_name} {last_name}".strip()
    mapping = {
        "CANDIDATE_NAME": full_name,
        "FIRST_NAME": first_name,
        "LAST_NAME": last_name,
        "POSITION": position,
        "SALARY": format_ars_dots(salary_num),
        "JOIN_DATE": fecha_es(join_date),
        "DATE": fecha_es(offer_date),
        "CITY": city,
    }
    extra_pairs = []
    for row in extras:
        k = (row.get("key") or "").strip()
        v = (row.get("value") or "").strip()
        if k:
            mapping[k.upper()] = v
            extra_pairs.append((k.upper(), v))

    try:
        edited = render_pptx(source_bytes, mapping)
        safe_name = " ".join(full_name.split()) or "Offer Letter"
        file_name_out = f"Offer Letter - {safe_name}.pptx"

        st.download_button(
            "Download Updated PPTX",
            data=edited.getvalue(),
            file_name=file_name_out,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        st.success(f"Done! Generated {file_name_out}.")

        # Build a thumbnail for history
        doc_thumb = pptx_doc_thumbnail(source_bytes)
        if doc_thumb:
            t_mime, t_bytes = doc_thumb
        else:
            # Placeholder using name + role
            t_mime, t_bytes = "image/svg+xml", svg_placeholder(full_name or "Offer", position or "")

        # Save in history
        push_history({
            "ts": datetime.now(),
            "file_name": file_name_out,
            "pptx_bytes": edited.getvalue(),
            "thumb_mime": t_mime,
            "thumb_bytes": t_bytes,
            "template": template_label,
            "fields": {
                "first_name": first_name, "last_name": last_name, "position": position,
                "salary_num": int(salary_num), "offer_date": offer_date.isoformat(),
                "join_date": join_date.isoformat(), "city": city,
            },
            "extras": extra_pairs,
        })
    except Exception as e:
        st.exception(e)

st.divider()

# ==============================
# History with thumbnails
# ==============================
ensure_history()
st.subheader(f"History (this session) — {len(st.session_state[HISTORY_KEY])} item(s)")

if st.session_state[HISTORY_KEY]:
    st.download_button(
        "Export all offers (ZIP)",
        data=zip_all_history(),
        file_name="offer_letters_history.zip",
        mime="application/zip",
        type="secondary",
    )

    # Cards newest-first
    for idx, item in reversed(list(enumerate(st.session_state[HISTORY_KEY]))):
        meta = item["fields"]
        ts = item["ts"].strftime("%Y-%m-%d %H:%M")

        cthumb, cmeta = st.columns([1, 3])
        with cthumb:
            if item["thumb_mime"] == "image/svg+xml":
                embed_svg(item["thumb_bytes"], width=220)
            elif item["thumb_mime"] in ("image/png", "image/jpeg"):
                st.image(item["thumb_bytes"], use_container_width=True)
            else:
                embed_svg(svg_placeholder("Offer Letter", f"{meta['first_name']} {meta['last_name']}"), width=220)

        with cmeta:
            st.markdown(
                f"**{item['file_name']}**  \n"
                f"**{meta['first_name']} {meta['last_name']}** — {meta['position']}  \n"
                f"{meta['city']} · Offer: {fecha_es(date.fromisoformat(meta['offer_date']))} · "
                f"Join: {fecha_es(date.fromisoformat(meta['join_date']))}  \n"
                f"Salary: {format_ars_dots(meta['salary_num'])}  \n"
                f"Template: {item.get('template','N/A')}  \n"
                f"*Created:* {ts}"
            )

            bcol1, bcol2, bcol3 = st.columns([1,1,1])
            with bcol1:
                st.download_button(
                    "Download again",
                    data=item["pptx_bytes"],
                    file_name=item["file_name"],
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key=f"dl_{idx}",
                )
            with bcol2:
                if st.button("Restore to form", key=f"restore_{idx}"):
                    st.session_state["First name"] = meta["first_name"]
                    st.session_state["Last name"]  = meta["last_name"]
                    st.session_state["Position"]   = meta["position"]
                    st.session_state["Salary (ARS)"] = meta["salary_num"]
                    st.session_state["Offer date"] = date.fromisoformat(meta["offer_date"])
                    st.session_state["Join date"]  = date.fromisoformat(meta["join_date"])
                    st.session_state["City"]       = meta["city"]
                    st.session_state["_extras_prefill"] = [{"key": k, "value": v} for k, v in item["extras"]]
                    st.experimental_rerun()
            with bcol3:
                if st.button("Delete", key=f"del_{idx}"):
                    delete_history(idx)
                    st.experimental_rerun()
else:
    st.info("No offers generated yet. When you create one, it will appear here with a thumbnail.")

# Consume extras prefill on restore (cannot live-update the editor mid-run)
if "_extras_prefill" in st.session_state:
    st.session_state.pop("_extras_prefill")

# ==============================
# Optional footer logo (add hogarth_split_black.png next to app.py)
# ==============================
def footer_logo():
    for name in ("hogarth_split_black.png", "hogarth_split.png"):
        p = Path(__file__).with_name(name)
        if p.exists():
            b64 = base64.b64encode(p.read_bytes()).decode("utf-8")
            st.markdown(
                f"<div style='display:flex;justify-content:center;margin-top:36px'>"
                f"<img src='data:image/png;base64,{b64}' style='max-width:520px;width:60%;height:auto' />"
                f"</div>",
                unsafe_allow_html=True,
            )
            break
footer_logo()
